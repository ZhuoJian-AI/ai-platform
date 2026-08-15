"""RAG / 工作空间共享文档解析器。

按扩展名 / Content-Type 分派到对应解析库（惰性 import，避免无解析任务时强依赖）：
- txt / md / csv：编码探测后解码
- html / htm：去标签取文本
- pdf：逐页抽取文本
- docx：标题、段落与表格
- xlsx：逐工作表、逐行单元格
- pptx：逐幻灯片标题、正文、表格与备注
- 旧版、模板、宏和 OpenDocument 变体：LibreOffice headless 转为现代格式后解析

调用方（``rag_service.ingest_uploaded_file``）在请求线程内同步调用本模块完成文本
抽取，结果立即落库为 ``RagDocument.content``；分块与嵌入在后台任务中异步进行。
"""

from __future__ import annotations

import io
import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from urllib.parse import unquote

import structlog

logger = structlog.get_logger()

# 单文件大小上限（50MB）——防止超大文件耗尽内存 / 拖垮解析
MAX_FILE_BYTES = 50 * 1024 * 1024
LIBREOFFICE_TIMEOUT_SECONDS = 30

WORD_CONVERT_EXTS = {"doc", "docm", "dot", "dotx", "dotm", "rtf", "odt"}
SHEET_CONVERT_EXTS = {"xls", "xlsb", "xlt", "xltx", "xltm", "ods"}
SLIDE_CONVERT_EXTS = {"ppt", "pptm", "pps", "ppsx", "ppsm", "pot", "potx", "potm", "odp"}

SUPPORTED_EXTENSIONS = {
    "txt", "text", "log", "md", "markdown", "csv", "tsv", "htm", "html", "pdf",
    "docx", "xlsx", "xlsm", "pptx",
    *WORD_CONVERT_EXTS, *SHEET_CONVERT_EXTS, *SLIDE_CONVERT_EXTS,
}


class UnsupportedFileTypeError(ValueError):
    """不支持的文件类型 / 无法解析。"""


def _ext(filename: str) -> str:
    name = unquote(filename or "").strip().lower()
    if "." not in name:
        return ""
    return name.rsplit(".", 1)[-1]


def _decode(raw: bytes) -> str:
    """字节流按探测编码解码，回退 utf-8（errors='replace' 不抛）。"""
    try:
        from charset_normalizer import from_bytes
        result = from_bytes(raw).best()
        if result is not None:
            return str(result)
    except Exception as exc:  # noqa: BLE001 — 探测失败回退 utf-8
        logger.debug("charset_detect_failed", error=str(exc))
    return raw.decode("utf-8", errors="replace")


def extract_text(filename: str, content_type: str | None, raw: bytes) -> tuple[str, str]:
    """抽取纯文本。返回 ``(text, kind)``，``kind`` 为解析类型标识（供 metadata）。

    无法解析时抛 :class:`UnsupportedFileTypeError`，由调用方转为 422 / 写入 parse_error。
    """
    if len(raw) > MAX_FILE_BYTES:
        raise UnsupportedFileTypeError(
            f"文件过大（{len(raw)} 字节），上限 {MAX_FILE_BYTES // (1024 * 1024)}MB"
        )
    if not raw:
        raise UnsupportedFileTypeError("文件为空")

    ext = _ext(filename)
    ct = (content_type or "").split(";")[0].strip().lower()

    kind = _resolve_kind(ext, ct)
    if kind is None:
        raise UnsupportedFileTypeError(f"不支持的文件类型：{ext or ct or '未知'}")

    try:
        if ext in WORD_CONVERT_EXTS:
            raw = _convert_with_libreoffice(filename, raw, "docx")
            kind = "docx"
        elif ext in SHEET_CONVERT_EXTS:
            raw = _convert_with_libreoffice(filename, raw, "xlsx")
            kind = "xlsx"
        elif ext in SLIDE_CONVERT_EXTS:
            raw = _convert_with_libreoffice(filename, raw, "pptx")
            kind = "pptx"
        if kind in ("txt", "md", "csv"):
            return _decode(raw), kind
        if kind == "html":
            return _parse_html(raw), kind
        if kind == "pdf":
            return _parse_pdf(raw), kind
        if kind == "docx":
            return _parse_docx(raw), kind
        if kind == "xlsx":
            return _parse_xlsx(raw), kind
        if kind == "pptx":
            return _parse_pptx(raw), kind
    except UnsupportedFileTypeError:
        raise
    except Exception as exc:  # noqa: BLE001 — 统一转为带原因的解析错误
        raise UnsupportedFileTypeError(f"{kind} 解析失败：{exc}") from exc

    # 理论不可达（kind 已归一化到上述分支）
    raise UnsupportedFileTypeError(f"不支持的文件类型：{ext or ct or '未知'}")


def _resolve_kind(ext: str, ct: str) -> str | None:
    """按扩展名（优先）或 Content-Type 归一化到解析类型。"""
    table = {
        "txt": "txt", "text": "txt", "log": "txt",
        "md": "md", "markdown": "md",
        "csv": "csv", "tsv": "csv",
        "htm": "html", "html": "html",
        "pdf": "pdf",
        "docx": "docx",
        "xlsx": "xlsx", "xlsm": "xlsx",
        "pptx": "pptx",
    }
    for item in WORD_CONVERT_EXTS:
        table[item] = "docx"
    for item in SHEET_CONVERT_EXTS:
        table[item] = "xlsx"
    for item in SLIDE_CONVERT_EXTS:
        table[item] = "pptx"
    if ext in table:
        return table[ext]
    # Content-Type 兜底
    ct_map = {
        "text/plain": "txt", "text/markdown": "md", "text/csv": "csv",
        "text/html": "html", "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    }
    return ct_map.get(ct)


def _convert_with_libreoffice(filename: str, raw: bytes, target_ext: str) -> bytes:
    """用独立 headless LibreOffice 进程把旧 Office 格式转换为现代中间格式。"""
    executable = shutil.which("soffice") or shutil.which("libreoffice")
    if executable is None:
        raise UnsupportedFileTypeError("服务器未安装 LibreOffice，无法解析旧版 Office 文件")

    suffix = f".{_ext(filename) or 'bin'}"
    with tempfile.TemporaryDirectory(prefix="ai-office-") as tmp:
        root = Path(tmp)
        source = root / f"source{suffix}"
        output = root / "output"
        profile = root / "profile"
        output.mkdir()
        profile.mkdir()
        source.write_bytes(raw)
        command = [
            executable,
            "--headless", "--safe-mode", "--nologo", "--nodefault",
            "--nofirststartwizard", "--norestore",
            f"-env:UserInstallation={profile.as_uri()}",
            "--convert-to", target_ext,
            "--outdir", str(output),
            str(source),
        ]
        env = {**os.environ, "HOME": str(root), "SAL_USE_VCLPLUGIN": "svp"}
        try:
            result = subprocess.run(
                command, capture_output=True, check=False, timeout=LIBREOFFICE_TIMEOUT_SECONDS,
                env=env,
            )
        except subprocess.TimeoutExpired as exc:
            raise UnsupportedFileTypeError("Office 文件转换超时") from exc
        converted = next(output.glob(f"*.{target_ext}"), None)
        if result.returncode != 0 or converted is None or not converted.is_file():
            detail = (result.stderr or result.stdout).decode("utf-8", errors="replace").strip()
            raise UnsupportedFileTypeError(f"Office 文件转换失败：{detail[:300] or '未知错误'}")
        return converted.read_bytes()


def _parse_html(raw: bytes) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_decode(raw), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    return "\n".join(line.strip() for line in text.splitlines() if line.strip())


def _parse_pdf(raw: bytes) -> str:
    import pdfplumber

    pages: list[str] = []
    with pdfplumber.open(io.BytesIO(raw)) as pdf:
        for index, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"# 第 {index} 页\n\n{text.strip()}")
    return "\n\n".join(pages)


def _parse_docx(raw: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(raw))
    parts: list[str] = []
    for para in document.paragraphs:
        text = (para.text or "").strip()
        if text:
            style = (getattr(para.style, "name", "") or "").lower()
            if style.startswith("heading"):
                level = style.removeprefix("heading").strip()
                hashes = "#" * max(1, min(int(level) if level.isdigit() else 1, 6))
                parts.append(f"{hashes} {text}")
            else:
                parts.append(text)
    # 表格单元格文本
    for table_index, table in enumerate(document.tables, start=1):
        rows: list[list[str]] = []
        for row in table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            if any(cells):
                rows.append(cells)
        if rows:
            parts.append(f"## 表格 {table_index}\n\n{_markdown_table(rows)}")
    return "\n\n".join(parts)


def _parse_xlsx(raw: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        rows: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            cells = [("" if v is None else str(v)).strip() for v in row]
            if any(cells):
                rows.append(cells)
        if rows:
            parts.append(f"# 工作表：{ws.title}\n\n{_markdown_table(rows)}")
    return "\n\n".join(parts)


def _parse_pptx(raw: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(raw))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = [f"# 幻灯片 {index}"]
        seen: set[str] = set()
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = (shape.text or "").strip()
                if text and text not in seen:
                    seen.add(text)
                    slide_parts.append(text)
            if getattr(shape, "has_table", False):
                rows = [[(cell.text or "").strip() for cell in row.cells] for row in shape.table.rows]
                if any(any(row) for row in rows):
                    slide_parts.append(_markdown_table(rows))
        try:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes and notes not in seen:
                slide_parts.append(f"## 备注\n\n{notes}")
        except (AttributeError, ValueError):
            pass
        parts.append("\n\n".join(slide_parts))
    return "\n\n".join(parts)


def _markdown_table(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    def escape(value: str) -> str:
        return value.replace("|", "\\|").replace("\n", "<br>")

    lines = ["| " + " | ".join(escape(value) for value in normalized[0]) + " |"]
    lines.append("| " + " | ".join("---" for _ in range(width)) + " |")
    lines.extend("| " + " | ".join(escape(value) for value in row) + " |" for row in normalized[1:])
    return "\n".join(lines)
