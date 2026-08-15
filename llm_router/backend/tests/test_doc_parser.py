"""Office/PDF 共享解析器的快速回归测试。"""

from __future__ import annotations

import io
import subprocess
from pathlib import Path
from types import SimpleNamespace

import pytest

from app.services import doc_parser


@pytest.fixture(autouse=True)
def db_engine():
    """纯解析单测不启动 conftest 的 PostgreSQL autouse fixture。"""
    yield


def _docx_bytes() -> bytes:
    from docx import Document

    output = io.BytesIO()
    document = Document()
    document.add_heading("生产日报", level=1)
    document.add_paragraph("今日完成 120 件")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "产品"
    table.cell(0, 1).text = "数量"
    table.cell(1, 0).text = "A 型"
    table.cell(1, 1).text = "120"
    document.save(output)
    return output.getvalue()


def _xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    output = io.BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "库存"
    sheet.append(["物料", "库存"])
    sheet.append(["螺丝", 300])
    workbook.save(output)
    return output.getvalue()


def _pptx_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    output = io.BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "销售月报"
    slide.placeholders[1].text = "本月增长 18%"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(6), Inches(1.5)).table
    table.cell(0, 0).text = "区域"
    table.cell(0, 1).text = "销售额"
    table.cell(1, 0).text = "华东"
    table.cell(1, 1).text = "100万"
    presentation.save(output)
    return output.getvalue()


def _pdf_bytes() -> bytes:
    """生成一个仅含文字层的最小 PDF，不额外引入 PDF 生成依赖。"""
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length 51 >>\nstream\nBT /F1 12 Tf 72 720 Td (Production report 42) Tj ET\nendstream",
    ]
    content = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, body in enumerate(objects, start=1):
        offsets.append(len(content))
        content.extend(f"{index} 0 obj\n".encode())
        content.extend(body)
        content.extend(b"\nendobj\n")
    xref_offset = len(content)
    content.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    content.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        content.extend(f"{offset:010d} 00000 n \n".encode())
    content.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF\n".encode()
    )
    return bytes(content)


@pytest.mark.parametrize(
    ("filename", "content_type", "raw_factory", "kind", "needles"),
    [
        (
            "生产日报.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            _docx_bytes,
            "docx",
            ["# 生产日报", "今日完成 120 件", "| 产品 | 数量 |"],
        ),
        (
            "库存.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            _xlsx_bytes,
            "xlsx",
            ["# 工作表：库存", "| 物料 | 库存 |", "| 螺丝 | 300 |"],
        ),
        (
            "销售月报.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            _pptx_bytes,
            "pptx",
            ["# 幻灯片 1", "销售月报", "本月增长 18%", "| 区域 | 销售额 |"],
        ),
    ],
)
def test_extract_modern_office(
    filename: str,
    content_type: str,
    raw_factory,
    kind: str,
    needles: list[str],
) -> None:
    text, actual_kind = doc_parser.extract_text(filename, content_type, raw_factory())

    assert actual_kind == kind
    for needle in needles:
        assert needle in text


def test_extract_pdf_text_by_page() -> None:
    text, kind = doc_parser.extract_text("质量报告.pdf", "application/pdf", _pdf_bytes())

    assert kind == "pdf"
    assert "# 第 1 页" in text
    assert "Production report 42" in text


def test_old_office_uses_isolated_libreoffice_profile(monkeypatch: pytest.MonkeyPatch) -> None:
    converted = _docx_bytes()
    observed: dict[str, object] = {}

    monkeypatch.setattr(doc_parser.shutil, "which", lambda _: "C:/LibreOffice/soffice.exe")

    def fake_run(command, *, capture_output, check, timeout, env):
        observed.update(command=command, timeout=timeout, env=env)
        outdir = Path(command[command.index("--outdir") + 1])
        (outdir / "source.docx").write_bytes(converted)
        return SimpleNamespace(returncode=0, stdout=b"", stderr=b"")

    monkeypatch.setattr(doc_parser.subprocess, "run", fake_run)

    text, kind = doc_parser.extract_text("历史报告.doc", "application/msword", b"legacy-doc")

    command = observed["command"]
    assert isinstance(command, list)
    assert "--headless" in command
    assert "--convert-to" in command
    assert command[command.index("--convert-to") + 1] == "docx"
    assert any(str(item).startswith("-env:UserInstallation=file:") for item in command)
    assert observed["timeout"] == doc_parser.LIBREOFFICE_TIMEOUT_SECONDS
    assert kind == "docx"
    assert "生产日报" in text


def test_old_office_conversion_timeout_is_clear(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(doc_parser.shutil, "which", lambda _: "soffice")

    def fake_timeout(*args, **kwargs):
        raise subprocess.TimeoutExpired(cmd="soffice", timeout=30)

    monkeypatch.setattr(doc_parser.subprocess, "run", fake_timeout)

    with pytest.raises(doc_parser.UnsupportedFileTypeError, match="转换超时"):
        doc_parser.extract_text("超时.xls", "application/vnd.ms-excel", b"legacy-xls")


@pytest.mark.parametrize(
    "filename",
    [
        "a.doc", "a.docx", "a.docm", "a.dot", "a.dotx", "a.dotm", "a.rtf", "a.odt",
        "a.xls", "a.xlsx", "a.xlsm", "a.xlsb", "a.xlt", "a.xltx", "a.xltm", "a.ods",
        "a.csv", "a.tsv", "a.ppt", "a.pptx", "a.pptm", "a.pps", "a.ppsx", "a.ppsm",
        "a.pot", "a.potx", "a.potm", "a.odp", "a.pdf",
    ],
)
def test_declared_business_formats_are_recognized(filename: str) -> None:
    assert doc_parser._resolve_kind(doc_parser._ext(filename), "application/octet-stream") is not None


def test_unsupported_extension_returns_clear_error() -> None:
    with pytest.raises(doc_parser.UnsupportedFileTypeError, match="不支持的文件类型"):
        doc_parser.extract_text("archive.rar", "application/octet-stream", b"not-a-rar")
