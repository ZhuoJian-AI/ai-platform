"""Deterministic platform-owned file tools executed inside skill-runner.

These handlers deliberately do not load a Skill package.  They use the
Runner's immutable base image while user Skills continue to execute through
their package-specific virtualenv/node_modules in ``app.py``.
"""

from __future__ import annotations

import csv
import html
import ipaddress
import json
import mimetypes
import re
import shutil
import socket
import stat
import subprocess
import tarfile
import tempfile
import urllib.parse
import urllib.request
import zipfile
from datetime import date, datetime, time
from decimal import Decimal
from html.parser import HTMLParser
from pathlib import Path, PurePosixPath
from typing import Any, ClassVar

import fitz
from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader, PdfWriter


class BuiltinToolError(ValueError):
    """A platform file operation could not be completed."""


MAX_WEB_TEXT_BYTES = 2 * 1024 * 1024
MAX_WEB_DOWNLOAD_BYTES = 5 * 1024 * 1024
MAX_ARCHIVE_FILES = 20
MAX_ARCHIVE_FILE_BYTES = 5 * 1024 * 1024
Image.MAX_IMAGE_PIXELS = 40_000_000


def _value(value: Any) -> Any:
    if isinstance(value, (str, int, float, bool)) or value is None:
        return value
    if isinstance(value, (datetime, date, time)):
        return value.isoformat()
    if isinstance(value, Decimal):
        return str(value)
    # Excel may contain library-specific scalar values. Inspection must never
    # turn an otherwise readable workbook into a Runner 500 response.
    return json.dumps(value, ensure_ascii=False, default=str)


def _safe_output_name(value: str | None, default: str, suffix: str) -> str:
    name = Path((value or default).replace("\\", "/")).name
    if not name.lower().endswith(suffix):
        name += suffix
    return name


def _libreoffice_convert(source: Path, output_dir: Path, target: str) -> Path:
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if not executable:
        raise BuiltinToolError("LibreOffice is unavailable in Runner")
    profile = Path(tempfile.mkdtemp(prefix="builtin-lo-"))
    try:
        result = subprocess.run(
            [
                executable,
                "--headless",
                f"-env:UserInstallation={profile.as_uri()}",
                "--convert-to",
                target,
                "--outdir",
                str(output_dir),
                str(source),
            ],
            capture_output=True,
            text=True,
            timeout=30,
            check=False,
        )
        if result.returncode:
            raise BuiltinToolError((result.stderr or result.stdout or "LibreOffice conversion failed")[-2000:])
        candidates = sorted(output_dir.glob(f"{source.stem}.*"), key=lambda item: item.stat().st_mtime)
        if not candidates:
            raise BuiltinToolError("LibreOffice did not produce an output file")
        return candidates[-1]
    except subprocess.TimeoutExpired as exc:
        raise BuiltinToolError("LibreOffice conversion exceeded 30 seconds") from exc
    finally:
        shutil.rmtree(profile, ignore_errors=True)


def _load_tabular(path: Path):
    if path.suffix.lower() in {".csv", ".tsv"}:
        delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
        rows = list(csv.reader(path.read_text(encoding="utf-8-sig").splitlines(), delimiter=delimiter))
        book = Workbook()
        sheet = book.active
        sheet.title = path.stem[:31] or "Sheet1"
        for row in rows:
            sheet.append(row)
        return book
    if path.suffix.lower() not in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
        converted_dir = Path(tempfile.mkdtemp(prefix="builtin-sheet-"))
        converted = _libreoffice_convert(path, converted_dir, "xlsx")
        book = load_workbook(converted, data_only=False)
        shutil.rmtree(converted_dir, ignore_errors=True)
        return book
    return load_workbook(path, data_only=False)


def _style_sheet(sheet) -> None:
    if sheet.max_row:
        for cell in sheet[1]:
            cell.font = Font(bold=True)
            cell.fill = PatternFill("solid", fgColor="DDEBF7")
            cell.alignment = Alignment(vertical="center")
    for column in range(1, min(sheet.max_column, 100) + 1):
        width = 10
        for row in range(1, min(sheet.max_row, 200) + 1):
            value = sheet.cell(row, column).value
            if value is not None:
                width = max(width, min(len(str(value)) + 2, 42))
        sheet.column_dimensions[get_column_letter(column)].width = width


def _write_sheets(book: Workbook, sheets: list[dict]) -> None:
    while book.worksheets:
        book.remove(book.worksheets[0])
    for index, spec in enumerate(sheets or [{"name": "Sheet1", "rows": []}]):
        sheet = book.create_sheet(str(spec.get("name") or f"Sheet{index + 1}")[:31])
        for row in spec.get("rows") or []:
            sheet.append([_value(value) for value in row])
        _style_sheet(sheet)


def _spreadsheet(action: str, inputs: list[Path], params: dict, output_dir: Path) -> dict:
    if action == "inspect":
        if not inputs:
            raise BuiltinToolError("spreadsheet inspect requires one input file")
        book = _load_tabular(inputs[0])
        max_rows = min(max(int(params.get("max_rows", 50)), 1), 200)
        max_cols = min(max(int(params.get("max_columns", 30)), 1), 100)
        sheets = []
        for sheet in book.worksheets:
            rows = [
                [_value(cell.value) for cell in row[:max_cols]]
                for row in sheet.iter_rows(min_row=1, max_row=min(sheet.max_row, max_rows))
            ]
            sheets.append({"name": sheet.title, "rows": rows, "total_rows": sheet.max_row,
                           "total_columns": sheet.max_column})
        return {"summary": {"kind": "spreadsheet", "sheets": sheets}, "outputs": []}

    if action == "create":
        book = Workbook()
        _write_sheets(book, params.get("sheets") or [{"name": "Sheet1", "rows": params.get("rows") or []}])
    elif action == "edit":
        if not inputs:
            raise BuiltinToolError("spreadsheet edit requires one input file")
        book = _load_tabular(inputs[0])
        for operation in params.get("operations") or []:
            op = operation.get("type")
            sheet_name = str(operation.get("sheet") or book.sheetnames[0])
            if sheet_name not in book.sheetnames:
                book.create_sheet(sheet_name[:31])
            sheet = book[sheet_name[:31]]
            if op == "set_cell":
                sheet[str(operation.get("cell") or "A1")] = _value(operation.get("value"))
            elif op == "append_rows":
                for row in operation.get("rows") or []:
                    sheet.append([_value(value) for value in row])
            elif op == "replace_rows":
                sheet.delete_rows(1, sheet.max_row)
                for row in operation.get("rows") or []:
                    sheet.append([_value(value) for value in row])
            elif op == "rename_sheet":
                sheet.title = str(operation.get("new_name") or sheet.title)[:31]
            else:
                raise BuiltinToolError(f"Unsupported spreadsheet edit operation: {op}")
        for sheet in book.worksheets:
            _style_sheet(sheet)
    elif action == "convert":
        if not inputs:
            raise BuiltinToolError("spreadsheet convert requires one input file")
        target = str(params.get("target_format") or "xlsx").lower().lstrip(".")
        if target not in {"xlsx", "csv", "tsv", "ods", "pdf"}:
            raise BuiltinToolError(f"Unsupported spreadsheet target format: {target}")
        if target in {"ods", "pdf"}:
            produced = _libreoffice_convert(inputs[0], output_dir, target)
            requested = _safe_output_name(params.get("output_name"), produced.name, f".{target}")
            final = output_dir / requested
            if produced != final:
                produced.replace(final)
            return {"summary": f"converted spreadsheet to {target}", "outputs": [final]}
        book = _load_tabular(inputs[0])
        if target in {"csv", "tsv"}:
            suffix = f".{target}"
            output = output_dir / _safe_output_name(params.get("output_name"), inputs[0].stem, suffix)
            delimiter = "\t" if target == "tsv" else ","
            with output.open("w", encoding="utf-8-sig", newline="") as handle:
                writer = csv.writer(handle, delimiter=delimiter)
                for row in book[book.sheetnames[0]].iter_rows(values_only=True):
                    writer.writerow(list(row))
            return {"summary": f"converted spreadsheet to {target}", "outputs": [output]}
    else:
        raise BuiltinToolError(f"Unsupported spreadsheet action: {action}")

    output = output_dir / _safe_output_name(params.get("output_name"), "workbook.xlsx", ".xlsx")
    book.save(output)
    return {"summary": f"{action}d spreadsheet", "outputs": [output]}


def _append_markdown(document: Document, markdown: str) -> None:
    lines = markdown.replace("\r\n", "\n").split("\n")
    index = 0
    while index < len(lines):
        line = lines[index]
        if line.startswith("### "):
            document.add_heading(line[4:], level=3)
        elif line.startswith("## "):
            document.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            document.add_heading(line[2:], level=1)
        elif line.startswith("- "):
            document.add_paragraph(line[2:], style="List Bullet")
        elif line[:3].rstrip(".").isdigit() and ". " in line[:5]:
            document.add_paragraph(line.split(". ", 1)[1], style="List Number")
        elif (
            line.startswith("|")
            and index + 1 < len(lines)
            and not lines[index + 1]
            .replace("|", "")
            .replace("-", "")
            .replace(":", "")
            .strip()
        ):
            rows: list[list[str]] = []
            while index < len(lines) and lines[index].startswith("|"):
                cells = [cell.strip() for cell in lines[index].strip("|").split("|")]
                if not all(set(cell) <= {"-", ":", " "} for cell in cells):
                    rows.append(cells)
                index += 1
            if rows:
                table = document.add_table(rows=len(rows), cols=max(len(row) for row in rows))
                table.style = "Table Grid"
                for row_index, row in enumerate(rows):
                    for column_index, value in enumerate(row):
                        table.cell(row_index, column_index).text = value
            continue
        elif line.strip():
            document.add_paragraph(line)
        index += 1


def _document(action: str, inputs: list[Path], params: dict, output_dir: Path) -> dict:
    if action == "inspect":
        if not inputs:
            raise BuiltinToolError("document inspect requires one input file")
        source = inputs[0]
        temp: Path | None = None
        if source.suffix.lower() != ".docx":
            temp = Path(tempfile.mkdtemp(prefix="builtin-doc-"))
            source = _libreoffice_convert(source, temp, "docx")
        document = Document(source)
        summary = {
            "kind": "document",
            "paragraphs": [paragraph.text for paragraph in document.paragraphs[:500]],
            "tables": [[[cell.text for cell in row.cells] for row in table.rows] for table in document.tables[:20]],
        }
        if temp:
            shutil.rmtree(temp, ignore_errors=True)
        return {"summary": summary, "outputs": []}
    if action == "create":
        document = Document()
        _append_markdown(document, str(params.get("markdown") or params.get("content") or ""))
    elif action == "edit":
        if not inputs:
            raise BuiltinToolError("document edit requires one DOCX input file")
        if inputs[0].suffix.lower() != ".docx":
            raise BuiltinToolError("document edit currently requires DOCX; convert legacy files first")
        if params.get("replace"):
            document = Document()
        else:
            document = Document(inputs[0])
        _append_markdown(document, str(params.get("markdown") or params.get("content") or ""))
    elif action == "convert":
        if not inputs:
            raise BuiltinToolError("document convert requires one input file")
        target = str(params.get("target_format") or "pdf").lower().lstrip(".")
        if target not in {"pdf", "docx", "odt", "rtf"}:
            raise BuiltinToolError(f"Unsupported document target format: {target}")
        produced = _libreoffice_convert(inputs[0], output_dir, target)
        final = output_dir / _safe_output_name(params.get("output_name"), produced.name, f".{target}")
        if produced != final:
            produced.replace(final)
        return {"summary": f"converted document to {target}", "outputs": [final]}
    else:
        raise BuiltinToolError(f"Unsupported document action: {action}")
    output = output_dir / _safe_output_name(params.get("output_name"), "document.docx", ".docx")
    document.save(output)
    return {"summary": f"{action}d document", "outputs": [output]}


def _add_slides(presentation: Presentation, slides: list[dict]) -> None:
    for spec in slides:
        layout = presentation.slide_layouts[1] if len(presentation.slide_layouts) > 1 else presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = str(spec.get("title") or "")
        body = "\n".join(str(value) for value in (spec.get("bullets") or []))
        placeholders = [shape for shape in slide.placeholders if shape != slide.shapes.title]
        if placeholders and hasattr(placeholders[0], "text_frame"):
            placeholders[0].text_frame.text = body
        elif body:
            box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
            box.text_frame.text = body
        if spec.get("notes") and slide.notes_slide.notes_text_frame:
            slide.notes_slide.notes_text_frame.text = str(spec["notes"])


def _presentation(action: str, inputs: list[Path], params: dict, output_dir: Path) -> dict:
    if action == "inspect":
        if not inputs:
            raise BuiltinToolError("presentation inspect requires one input file")
        source = inputs[0]
        temp: Path | None = None
        if source.suffix.lower() != ".pptx":
            temp = Path(tempfile.mkdtemp(prefix="builtin-ppt-"))
            source = _libreoffice_convert(source, temp, "pptx")
        presentation = Presentation(source)
        slides = []
        for index, slide in enumerate(presentation.slides):
            slides.append(
                {
                    "number": index + 1,
                    "texts": [
                        shape.text
                        for shape in slide.shapes
                        if hasattr(shape, "text") and shape.text
                    ],
                }
            )
        if temp:
            shutil.rmtree(temp, ignore_errors=True)
        return {"summary": {"kind": "presentation", "slides": slides}, "outputs": []}
    if action == "create":
        presentation = Presentation()
        _add_slides(presentation, params.get("slides") or [])
    elif action == "edit":
        if not inputs or inputs[0].suffix.lower() != ".pptx":
            raise BuiltinToolError("presentation edit requires one PPTX input file")
        presentation = Presentation(inputs[0])
        _add_slides(presentation, params.get("slides") or [])
    elif action == "convert":
        if not inputs:
            raise BuiltinToolError("presentation convert requires one input file")
        target = str(params.get("target_format") or "pdf").lower().lstrip(".")
        if target not in {"pdf", "pptx", "odp"}:
            raise BuiltinToolError(f"Unsupported presentation target format: {target}")
        produced = _libreoffice_convert(inputs[0], output_dir, target)
        final = output_dir / _safe_output_name(params.get("output_name"), produced.name, f".{target}")
        if produced != final:
            produced.replace(final)
        return {"summary": f"converted presentation to {target}", "outputs": [final]}
    else:
        raise BuiltinToolError(f"Unsupported presentation action: {action}")
    output = output_dir / _safe_output_name(params.get("output_name"), "presentation.pptx", ".pptx")
    presentation.save(output)
    return {"summary": f"{action}d presentation", "outputs": [output]}


def _pdf(action: str, inputs: list[Path], params: dict, output_dir: Path) -> dict:
    if action == "inspect":
        if not inputs:
            raise BuiltinToolError("pdf inspect requires one input file")
        reader = PdfReader(inputs[0])
        max_pages = min(max(int(params.get("max_pages", 20)), 1), 100)
        return {"summary": {"kind": "pdf", "page_count": len(reader.pages),
                            "pages": [(page.extract_text() or "")[:20_000] for page in reader.pages[:max_pages]]},
                "outputs": []}
    if action == "create":
        document = Document()
        _append_markdown(document, str(params.get("markdown") or params.get("content") or ""))
        temp_docx = output_dir / "source.docx"
        document.save(temp_docx)
        produced = _libreoffice_convert(temp_docx, output_dir, "pdf")
        temp_docx.unlink(missing_ok=True)
        final = output_dir / _safe_output_name(params.get("output_name"), "document.pdf", ".pdf")
        if produced != final:
            produced.replace(final)
        return {"summary": "created PDF", "outputs": [final]}
    if action == "edit":
        operation = str(params.get("operation") or "merge")
        writer = PdfWriter()
        if operation == "merge":
            if not inputs:
                raise BuiltinToolError("PDF merge requires input files")
            for source in inputs:
                for page in PdfReader(source).pages:
                    writer.add_page(page)
        elif operation in {"split", "extract_pages"}:
            if not inputs:
                raise BuiltinToolError("PDF page extraction requires one input file")
            reader = PdfReader(inputs[0])
            pages = params.get("pages") or [1]
            for number in pages:
                index = int(number) - 1
                if index < 0 or index >= len(reader.pages):
                    raise BuiltinToolError(f"PDF page {number} is out of range")
                writer.add_page(reader.pages[index])
        else:
            raise BuiltinToolError(f"Unsupported PDF edit operation: {operation}")
        output = output_dir / _safe_output_name(params.get("output_name"), "output.pdf", ".pdf")
        with output.open("wb") as handle:
            writer.write(handle)
        return {"summary": f"PDF {operation} completed", "outputs": [output]}
    if action == "convert":
        if not inputs:
            raise BuiltinToolError("pdf convert requires one input file")
        target = str(params.get("target_format") or "txt").lower().lstrip(".")
        if target != "txt":
            raise BuiltinToolError("PDF conversion currently supports txt output")
        text = "\n\n".join(page.extract_text() or "" for page in PdfReader(inputs[0]).pages)
        output = output_dir / _safe_output_name(params.get("output_name"), inputs[0].stem, ".txt")
        output.write_text(text, encoding="utf-8")
        return {"summary": "converted PDF to text", "outputs": [output]}
    raise BuiltinToolError(f"Unsupported PDF action: {action}")


def _text(action: str, inputs: list[Path], params: dict, output_dir: Path) -> dict:
    if action == "inspect":
        if not inputs:
            raise BuiltinToolError("text inspect requires one input file")
        content = inputs[0].read_text(encoding="utf-8-sig", errors="replace")
        return {"summary": {"kind": "text", "characters": len(content), "content": content[:100_000]}, "outputs": []}
    suffix = ".md" if str(params.get("format") or "").lower() in {"md", "markdown"} else ".txt"
    if action == "create":
        content = str(params.get("content") or "")
    elif action == "edit":
        if not inputs:
            raise BuiltinToolError("text edit requires one input file")
        original = inputs[0].read_text(encoding="utf-8-sig", errors="replace")
        if params.get("replace"):
            content = str(params.get("content") or "")
        else:
            content = original + str(params.get("content") or "")
        suffix = inputs[0].suffix.lower() if inputs[0].suffix.lower() in {".txt", ".md"} else suffix
    elif action == "convert":
        if not inputs:
            raise BuiltinToolError("text convert requires one input file")
        content = inputs[0].read_text(encoding="utf-8-sig", errors="replace")
        target = str(params.get("target_format") or "txt").lower().lstrip(".")
        if target not in {"txt", "md"}:
            raise BuiltinToolError(f"Unsupported text target format: {target}")
        suffix = f".{target}"
    else:
        raise BuiltinToolError(f"Unsupported text action: {action}")
    output = output_dir / _safe_output_name(params.get("output_name"), "document" + suffix, suffix)
    output.write_text(content, encoding="utf-8")
    return {"summary": f"{action}d text file", "outputs": [output]}


class _ReadableHTMLParser(HTMLParser):
    """Extract a useful title and readable text without running page scripts."""

    _SKIP_TAGS: ClassVar[set[str]] = {"script", "style", "noscript", "svg"}
    _BLOCK_TAGS: ClassVar[set[str]] = {
        "article", "aside", "blockquote", "br", "div", "footer", "h1", "h2",
        "h3", "h4", "h5", "h6", "header", "li", "main", "nav", "p", "section",
        "table", "td", "th", "tr",
    }

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._skip_depth = 0
        self._in_title = False
        self.title_parts: list[str] = []
        self.text_parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        del attrs
        tag = tag.lower()
        if tag in self._SKIP_TAGS:
            self._skip_depth += 1
        if self._skip_depth:
            return
        if tag == "title":
            self._in_title = True
        if tag in self._BLOCK_TAGS:
            self.text_parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if tag in self._SKIP_TAGS and self._skip_depth:
            self._skip_depth -= 1
            return
        if self._skip_depth:
            return
        if tag == "title":
            self._in_title = False
        if tag in self._BLOCK_TAGS:
            self.text_parts.append("\n")

    def handle_data(self, data: str) -> None:
        if self._skip_depth:
            return
        value = data.strip()
        if not value:
            return
        if self._in_title:
            self.title_parts.append(value)
        self.text_parts.append(value + " ")

    def result(self) -> tuple[str, str]:
        title = " ".join(self.title_parts).strip()
        text = "".join(self.text_parts)
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n\s*\n+", "\n\n", text)
        return title, text.strip()


class _DuckDuckGoParser(HTMLParser):
    """Parse the stable DuckDuckGo HTML result page without a browser."""

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.results: list[dict[str, str]] = []
        self._capture: str | None = None
        self._parts: list[str] = []
        self._href = ""

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        values = {key: value or "" for key, value in attrs}
        classes = set(values.get("class", "").split())
        if tag == "a" and "result__a" in classes:
            self._capture = "title"
            self._parts = []
            self._href = values.get("href", "")
        elif tag in {"a", "div"} and "result__snippet" in classes:
            self._capture = "snippet"
            self._parts = []

    def handle_endtag(self, tag: str) -> None:
        if self._capture == "title" and tag == "a":
            title = " ".join(self._parts).strip()
            if title and self._href:
                parsed = urllib.parse.urlparse(self._href)
                query = urllib.parse.parse_qs(parsed.query)
                href = query.get("uddg", [self._href])[0]
                self.results.append({"title": title, "url": href, "snippet": ""})
            self._capture = None
        elif self._capture == "snippet" and tag in {"a", "div"}:
            if self.results:
                self.results[-1]["snippet"] = " ".join(self._parts).strip()
            self._capture = None

    def handle_data(self, data: str) -> None:
        if self._capture and data.strip():
            self._parts.append(data.strip())


class _BingParser(HTMLParser):
    """Parse Bing's non-JavaScript result markup as a search fallback."""

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.results: list[dict[str, str]] = []
        self._current: dict[str, str] | None = None
        self._in_heading = False
        self._capture: str | None = None
        self._parts: list[str] = []

    def _finish_result(self) -> None:
        if self._current and self._current.get("title") and self._current.get("url"):
            self.results.append(self._current)
        self._current = None
        self._in_heading = False
        self._capture = None
        self._parts = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        values = {key: value or "" for key, value in attrs}
        classes = set(values.get("class", "").split())
        if tag == "li" and "b_algo" in classes:
            self._finish_result()
            self._current = {"title": "", "url": "", "snippet": ""}
            return
        if self._current is None:
            return
        if tag == "h2":
            self._in_heading = True
        elif tag == "a" and self._in_heading:
            self._capture = "title"
            self._parts = []
            self._current["url"] = values.get("href", "")
        elif tag == "p":
            self._capture = "snippet"
            self._parts = []

    def handle_endtag(self, tag: str) -> None:
        if self._current is None:
            return
        if tag == "a" and self._capture == "title":
            self._current["title"] = " ".join(self._parts).strip()
            self._capture = None
        elif tag == "p" and self._capture == "snippet":
            self._current["snippet"] = " ".join(self._parts).strip()
            self._capture = None
        elif tag == "h2":
            self._in_heading = False
        elif tag == "li":
            self._finish_result()

    def handle_data(self, data: str) -> None:
        if self._capture and data.strip():
            self._parts.append(data.strip())

    def close(self) -> None:
        super().close()
        self._finish_result()


def _validate_public_url(url: str) -> urllib.parse.ParseResult:
    parsed = urllib.parse.urlparse(url)
    if parsed.scheme not in {"http", "https"} or not parsed.hostname:
        raise BuiltinToolError("web_tool only accepts absolute HTTP/HTTPS URLs")
    hostname = parsed.hostname.rstrip(".").lower()
    if hostname == "localhost" or hostname.endswith(".local"):
        raise BuiltinToolError("Local or private network URLs are not allowed")
    try:
        port = parsed.port or (443 if parsed.scheme == "https" else 80)
        addresses = {item[4][0] for item in socket.getaddrinfo(hostname, port)}
    except ValueError as exc:
        raise BuiltinToolError("Invalid URL port") from exc
    except socket.gaierror as exc:
        raise BuiltinToolError(f"Unable to resolve host: {hostname}") from exc
    for value in addresses:
        address = ipaddress.ip_address(value.split("%", 1)[0])
        if (
            address.is_private
            or address.is_loopback
            or address.is_link_local
            or address.is_multicast
            or address.is_reserved
            or address.is_unspecified
        ):
            raise BuiltinToolError("Local or private network URLs are not allowed")
    return parsed


class _SafeRedirectHandler(urllib.request.HTTPRedirectHandler):
    def redirect_request(self, req, fp, code, msg, headers, newurl):
        _validate_public_url(newurl)
        return super().redirect_request(req, fp, code, msg, headers, newurl)


def _http_get(url: str, *, limit: int) -> tuple[str, bytes, str, str]:
    _validate_public_url(url)
    request = urllib.request.Request(
        url,
        headers={
            "User-Agent": "AI-Platform-WebTool/1.0 (+https://ai-platform.staging.zhuojianai.com)",
            "Accept": "text/html,application/json,text/plain,*/*;q=0.5",
        },
    )
    opener = urllib.request.build_opener(_SafeRedirectHandler())
    try:
        with opener.open(request, timeout=20) as response:
            raw = response.read(limit + 1)
            if len(raw) > limit:
                raise BuiltinToolError(f"Remote response exceeds {limit // (1024 * 1024)}MB")
            content_type = response.headers.get_content_type() or "application/octet-stream"
            disposition = response.headers.get("Content-Disposition", "")
            return response.geturl(), raw, content_type, disposition
    except BuiltinToolError:
        raise
    except Exception as exc:
        raise BuiltinToolError(f"Web request failed: {exc}") from exc


def _decode_web_text(raw: bytes, content_type: str) -> str:
    del content_type
    for encoding in ("utf-8", "gb18030", "big5"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _download_name(url: str, disposition: str, content_type: str) -> str:
    match = re.search(
        r"filename\*?=(?:UTF-8''|\")?([^\";]+)", disposition, flags=re.IGNORECASE
    )
    if match:
        name = urllib.parse.unquote(match.group(1).strip())
    else:
        name = urllib.parse.unquote(PurePosixPath(urllib.parse.urlparse(url).path).name)
    if not name:
        name = "download" + (mimetypes.guess_extension(content_type) or ".bin")
    return Path(name.replace("\\", "/")).name


def _web(action: str, inputs: list[Path], params: dict, output_dir: Path) -> dict:
    del inputs
    if action == "search":
        query = str(params.get("query") or "").strip()
        if not query:
            raise BuiltinToolError("web search requires query")
        max_results = min(max(int(params.get("max_results", 5)), 1), 10)
        url = "https://html.duckduckgo.com/html/?" + urllib.parse.urlencode({"q": query})
        _, raw, content_type, _ = _http_get(url, limit=MAX_WEB_TEXT_BYTES)
        parser = _DuckDuckGoParser()
        parser.feed(_decode_web_text(raw, content_type))
        results = parser.results[:max_results]
        if not results:
            fallback_url = "https://www.bing.com/search?" + urllib.parse.urlencode(
                {"q": query, "count": max_results}
            )
            _, raw, content_type, _ = _http_get(fallback_url, limit=MAX_WEB_TEXT_BYTES)
            fallback = _BingParser()
            fallback.feed(_decode_web_text(raw, content_type))
            fallback.close()
            results = fallback.results[:max_results]
        if not results:
            raise BuiltinToolError("Search providers returned no parseable results")
        return {"summary": {"query": query, "results": results}, "outputs": []}
    if action == "fetch":
        url = str(params.get("url") or "").strip()
        final_url, raw, content_type, _ = _http_get(url, limit=MAX_WEB_TEXT_BYTES)
        text = _decode_web_text(raw, content_type)
        title = ""
        if content_type in {"text/html", "application/xhtml+xml"} or "<html" in text[:500].lower():
            parser = _ReadableHTMLParser()
            parser.feed(text)
            title, text = parser.result()
        max_chars = min(max(int(params.get("max_chars", 50_000)), 1_000), 100_000)
        return {
            "summary": {
                "url": final_url,
                "title": html.unescape(title),
                "content_type": content_type,
                "content": text[:max_chars],
                "truncated": len(text) > max_chars,
            },
            "outputs": [],
        }
    if action == "download":
        url = str(params.get("url") or "").strip()
        final_url, raw, content_type, disposition = _http_get(
            url,
            limit=MAX_WEB_DOWNLOAD_BYTES,
        )
        name = _safe_output_name(
            params.get("output_name"),
            _download_name(final_url, disposition, content_type),
            "",
        )
        output = output_dir / name
        output.write_bytes(raw)
        return {
            "summary": {"url": final_url, "content_type": content_type, "size": len(raw)},
            "outputs": [output],
        }
    raise BuiltinToolError(f"Unsupported web action: {action}")


def _image_output_suffix(value: str | None, fallback: str = "png") -> tuple[str, str]:
    normalized = str(value or fallback).lower().lstrip(".")
    aliases = {"jpg": "jpeg", "tif": "tiff"}
    normalized = aliases.get(normalized, normalized)
    if normalized not in {"png", "jpeg", "webp", "tiff", "bmp"}:
        raise BuiltinToolError(f"Unsupported image format: {normalized}")
    suffix = ".jpg" if normalized == "jpeg" else f".{normalized}"
    return normalized.upper(), suffix


def _save_image(image: Image.Image, output: Path, image_format: str, quality: int) -> None:
    if image_format == "JPEG" and image.mode not in {"RGB", "L"}:
        background = Image.new("RGB", image.size, "white")
        if "A" in image.getbands():
            background.paste(image, mask=image.getchannel("A"))
        else:
            background.paste(image.convert("RGB"))
        image = background
    kwargs: dict[str, Any] = {"format": image_format}
    if image_format in {"JPEG", "WEBP"}:
        kwargs.update({"quality": quality, "optimize": True})
    elif image_format == "PNG":
        kwargs["optimize"] = True
    image.save(output, **kwargs)


def _ocr_images(source: Path, max_pages: int) -> list[Image.Image]:
    if source.suffix.lower() == ".pdf":
        document = fitz.open(source)
        images: list[Image.Image] = []
        try:
            for page in document[:max_pages]:
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                images.append(Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples))
        finally:
            document.close()
        return images
    with Image.open(source) as image:
        return [ImageOps.exif_transpose(image).convert("RGB")]


def _image(action: str, inputs: list[Path], params: dict, output_dir: Path) -> dict:
    if not inputs:
        raise BuiltinToolError(f"image {action} requires an input file")
    source = inputs[0]
    if action == "inspect":
        if source.suffix.lower() == ".pdf":
            document = fitz.open(source)
            try:
                return {
                    "summary": {"kind": "scanned_document", "pages": document.page_count},
                    "outputs": [],
                }
            finally:
                document.close()
        with Image.open(source) as image:
            return {
                "summary": {
                    "kind": "image",
                    "format": image.format,
                    "width": image.width,
                    "height": image.height,
                    "mode": image.mode,
                    "frames": getattr(image, "n_frames", 1),
                },
                "outputs": [],
            }
    if action == "ocr":
        try:
            import pytesseract
        except ImportError as exc:  # pragma: no cover - production image includes it
            raise BuiltinToolError("OCR runtime is unavailable") from exc
        language = str(params.get("language") or "chi_sim+eng")
        max_pages = min(max(int(params.get("max_pages", 10)), 1), 20)
        texts = [pytesseract.image_to_string(image, lang=language) for image in _ocr_images(source, max_pages)]
        content = "\n\n".join(text.strip() for text in texts).strip()
        outputs: list[Path] = []
        if params.get("output_name"):
            output = output_dir / _safe_output_name(params.get("output_name"), "ocr.txt", ".txt")
            output.write_text(content, encoding="utf-8")
            outputs.append(output)
        return {
            "summary": {"language": language, "pages": len(texts), "content": content[:100_000]},
            "outputs": outputs,
        }
    if source.suffix.lower() == ".pdf":
        raise BuiltinToolError(f"image {action} does not accept PDF; use ocr for scanned PDFs")
    with Image.open(source) as opened:
        image = ImageOps.exif_transpose(opened).copy()
    image_format, suffix = _image_output_suffix(
        params.get("target_format"),
        (source.suffix.lower().lstrip(".") or "png"),
    )
    if action == "resize":
        width = int(params.get("width") or image.width)
        height = int(params.get("height") or image.height)
        if width < 1 or height < 1 or width * height > 40_000_000:
            raise BuiltinToolError("Invalid or oversized image dimensions")
        if params.get("keep_aspect", True):
            image.thumbnail((width, height), Image.Resampling.LANCZOS)
        else:
            image = image.resize((width, height), Image.Resampling.LANCZOS)
    elif action == "crop":
        box = params.get("box") or []
        if not isinstance(box, list) or len(box) != 4:
            raise BuiltinToolError("image crop requires box=[left, top, right, bottom]")
        coordinates = tuple(int(value) for value in box)
        if coordinates[0] < 0 or coordinates[1] < 0 or coordinates[2] > image.width or coordinates[3] > image.height:
            raise BuiltinToolError("Crop box is outside the image")
        if coordinates[2] <= coordinates[0] or coordinates[3] <= coordinates[1]:
            raise BuiltinToolError("Crop box has no area")
        image = image.crop(coordinates)
    elif action not in {"convert", "compress"}:
        raise BuiltinToolError(f"Unsupported image action: {action}")
    quality = min(max(int(params.get("quality", 85)), 1), 100)
    output = output_dir / _safe_output_name(params.get("output_name"), "image" + suffix, suffix)
    _save_image(image, output, image_format, quality)
    return {
        "summary": {"action": action, "width": image.width, "height": image.height},
        "outputs": [output],
    }


def _safe_archive_path(name: str) -> PurePosixPath:
    normalized = PurePosixPath(name.replace("\\", "/").lstrip("/"))
    if not normalized.parts or ".." in normalized.parts or normalized.is_absolute():
        raise BuiltinToolError(f"Unsafe archive path: {name}")
    return normalized


def _archive_kind(path: Path) -> str:
    name = path.name.lower()
    if name.endswith(".zip"):
        return "zip"
    if name.endswith((".tar.gz", ".tgz")):
        return "tar.gz"
    if name.endswith(".tar"):
        return "tar"
    raise BuiltinToolError("archive_tool supports ZIP, TAR and TAR.GZ files")


def _archive_list(source: Path) -> list[dict[str, Any]]:
    kind = _archive_kind(source)
    if kind == "zip":
        with zipfile.ZipFile(source) as archive:
            return [
                {"path": item.filename, "size": item.file_size, "is_dir": item.is_dir()}
                for item in archive.infolist()
            ]
    with tarfile.open(source, "r:gz" if kind == "tar.gz" else "r:") as archive:
        return [
            {"path": item.name, "size": item.size, "is_dir": item.isdir()}
            for item in archive.getmembers()
        ]


def _archive_extract(source: Path, output_dir: Path) -> list[Path]:
    kind = _archive_kind(source)
    outputs: list[Path] = []
    if kind == "zip":
        with zipfile.ZipFile(source) as archive:
            members = [item for item in archive.infolist() if not item.is_dir()]
            if len(members) > MAX_ARCHIVE_FILES:
                raise BuiltinToolError(f"Archive contains more than {MAX_ARCHIVE_FILES} files")
            for item in members:
                if item.flag_bits & 0x1:
                    raise BuiltinToolError("Encrypted ZIP files are not supported")
                if stat.S_IFMT(item.external_attr >> 16) == stat.S_IFLNK:
                    raise BuiltinToolError("Archive links are not supported")
                if item.file_size > MAX_ARCHIVE_FILE_BYTES:
                    raise BuiltinToolError(f"Archive member {item.filename} exceeds 5MB")
                relative = _safe_archive_path(item.filename)
                target = output_dir.joinpath(*relative.parts)
                target.parent.mkdir(parents=True, exist_ok=True)
                with archive.open(item) as source_file, target.open("wb") as target_file:
                    shutil.copyfileobj(source_file, target_file)
                outputs.append(target)
        return outputs
    with tarfile.open(source, "r:gz" if kind == "tar.gz" else "r:") as archive:
        members = [item for item in archive.getmembers() if item.isfile()]
        if len(members) > MAX_ARCHIVE_FILES:
            raise BuiltinToolError(f"Archive contains more than {MAX_ARCHIVE_FILES} files")
        for item in members:
            if item.issym() or item.islnk():
                raise BuiltinToolError("Archive links are not supported")
            if item.size > MAX_ARCHIVE_FILE_BYTES:
                raise BuiltinToolError(f"Archive member {item.name} exceeds 5MB")
            relative = _safe_archive_path(item.name)
            target = output_dir.joinpath(*relative.parts)
            target.parent.mkdir(parents=True, exist_ok=True)
            source_file = archive.extractfile(item)
            if source_file is None:
                continue
            with source_file, target.open("wb") as target_file:
                shutil.copyfileobj(source_file, target_file)
            outputs.append(target)
    return outputs


def _archive(action: str, inputs: list[Path], params: dict, output_dir: Path) -> dict:
    if action == "list":
        if len(inputs) != 1:
            raise BuiltinToolError("archive list requires one input file")
        members = _archive_list(inputs[0])
        return {
            "summary": {
                "format": _archive_kind(inputs[0]),
                "total": len(members),
                "members": members[:200],
                "truncated": len(members) > 200,
            },
            "outputs": [],
        }
    if action == "extract":
        if len(inputs) != 1:
            raise BuiltinToolError("archive extract requires one input file")
        outputs = _archive_extract(inputs[0], output_dir)
        return {"summary": {"extracted": len(outputs)}, "outputs": outputs}
    if action == "create":
        if not inputs:
            raise BuiltinToolError("archive create requires input files")
        kind = str(params.get("format") or "zip").lower()
        if kind not in {"zip", "tar", "tar.gz", "tgz"}:
            raise BuiltinToolError("archive format must be zip, tar or tar.gz")
        suffix = ".zip" if kind == "zip" else (".tar" if kind == "tar" else ".tar.gz")
        output = output_dir / _safe_output_name(params.get("output_name"), "archive" + suffix, suffix)
        if kind == "zip":
            with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
                for source in inputs:
                    archive.write(source, arcname=source.name)
        else:
            with tarfile.open(output, "w" if kind == "tar" else "w:gz") as archive:
                for source in inputs:
                    archive.add(source, arcname=source.name, recursive=False)
        return {"summary": {"created": output.name, "files": len(inputs)}, "outputs": [output]}
    raise BuiltinToolError(f"Unsupported archive action: {action}")


def execute_builtin(tool_kind: str, action: str, inputs: list[Path], params: dict, output_dir: Path) -> dict:
    handlers = {
        "spreadsheet": _spreadsheet,
        "document": _document,
        "presentation": _presentation,
        "pdf": _pdf,
        "text": _text,
        "web": _web,
        "image": _image,
        "archive": _archive,
    }
    handler = handlers.get(tool_kind)
    if handler is None:
        raise BuiltinToolError(f"Unsupported builtin tool kind: {tool_kind}")
    result = handler(action, inputs, params, output_dir)
    result["mime_types"] = {
        path.name: mimetypes.guess_type(path.name)[0] or "application/octet-stream"
        for path in result.get("outputs") or []
    }
    return result
